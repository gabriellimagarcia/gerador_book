# === PARTE 1/10 =====================================================
# Boot, imports, logging e estilos iniciais

import re
import base64
import hashlib
from io import BytesIO
from collections import OrderedDict
from concurrent.futures import ThreadPoolExecutor, as_completed
import zipfile
import numpy as np
import logging, sys  # <<< LOGGING

import streamlit as st
import pandas as pd
import requests
from PIL import Image, ImageOps, ImageDraw, ImageFilter, ImageFont, ImageFile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 👇 ajustes de robustez do Pillow
ImageFile.LOAD_TRUNCATED_IMAGES = True      # aceita imagens truncadas
Image.MAX_IMAGE_PIXELS = 60_000_000         # evita DecompressionBombError para fotos muito grandes

# -------------------------------------------------------------------
# LOGGING (stdout -> aparece nos logs do Streamlit Cloud)
# -------------------------------------------------------------------
logging.basicConfig(
    level=logging.INFO,  # mude para logging.DEBUG para ver tudo
    format="%(asctime)s | %(levelname)s | %(name)s | %(message)s",
    stream=sys.stdout,
)
logger = logging.getLogger("gerador_book")

# -------------------------------------------------------------------
# CONFIG GERAL
# -------------------------------------------------------------------
st.set_page_config(page_title="Gerador de Book", page_icon="📸", layout="wide")
st.set_option("client.showErrorDetails", True)  # detalha erros na UI

# --- CSS básico (UX) ---
BASE_CSS = """
<style>
.steps {display:flex; gap:8px; align-items:center; margin:.25rem 0 .5rem;}
.step {padding:6px 10px; border-radius:999px; border:1px solid #E5E7EB; color:#111; background:#fff; font-weight:700; font-size:13px;}
.step.active {background:#FF7A00; color:#fff; border-color:#FF7A00;}
.step.sep {opacity:.5}

.img-card { border:1px solid #DDD; border-radius:10px; padding:10px; transition:all .15s ease;
  background:#fff; display:flex; flex-direction:column; gap:8px; height:100%; }
.dark .img-card { background:#0f131a; border-color:#232a36; }
.img-card:hover { box-shadow:0 8px 20px rgba(0,0,0,.06); transform:translateY(-2px); }
</style>
"""
st.markdown(BASE_CSS, unsafe_allow_html=True)


# === PARTE 2/10 =====================================================
# Continuação dos estilos (CSS login + tema)

LOGIN_CSS = """
<style>
.login-wrap { max-width: 520px; margin: 0 auto; }
.login-hero {
  background: linear-gradient(90deg, #FF7A00 0%, #FF9944 100%);
  color:#fff; padding:16px 18px; border-radius:14px 14px 0 0;
  font-weight:700; line-height:1.2; margin: 8px 0 0 0;
}
.login-card {
  border:1px solid #E7E7E7; border-top:none; background:#ffffff;
  padding:18px; border-radius:0 0 14px 14px; margin:0;
}
.stForm .stButton > button {
  background:#FF7A00 !important; color:#fff !important;
  border:none !important; border-radius:10px !important;
  font-weight:800 !important;
}
.stForm .stButton > button:hover { background:#E66E00 !important; }
</style>
"""

# -------------------------------------------------------------------
# TEMA
# -------------------------------------------------------------------
def apply_theme(dark: bool):
    ORANGE = "#FF7A00"; ORANGE_HOVER = "#E66E00"; BLACK = "#111111"; GRAY_BG = "#f6f6f7"
    if dark:
        palette = f"""
        <style>
        :root {{ --accent:{ORANGE}; --accent-hover:{ORANGE_HOVER}; --text:#f5f5f5; --bg:#0e1117; --panel:#11151c; }}
        .stApp {{ background-color:var(--bg); color:var(--text); }}
        section[data-testid="stSidebar"] > div {{ background:var(--panel); border-right:1px solid #1b212c; }}
        .stButton > button, .stDownloadButton > button {{ background:var(--accent); color:#fff; border-radius:10px; border:none; }}
        .stButton > button:hover, .stDownloadButton > button:hover {{ background:var(--accent-hover); }}
        .stProgress > div > div {{ background-color:var(--accent); }}
        </style>
        """
    else:
        palette = f"""
        <style>
        :root {{ --accent:{ORANGE}; --accent-hover:{ORANGE_HOVER}; --text:{BLACK}; --bg:#ffffff; --panel:{GRAY_BG}; }}
        .stApp {{ background-color:var(--bg); color:var(--text); }}
        section[data-testid="stSidebar"] > div {{ background:var(--panel); border-right:1px solid #ececec; }}
        .stButton > button, .stDownloadButton > button {{ background:var(--accent); color:#fff; border-radius:10px; border:none; }}
        .stButton > button:hover, .stDownloadButton > button:hover {{ background:var(--accent-hover); }}
        .stProgress > div > div {{ background-color:var(--accent); }}
        </style>
        """
    st.markdown(palette, unsafe_allow_html=True)

if "dark_mode" not in st.session_state:
    st.session_state.dark_mode = False


# === PARTE 3/10 =====================================================
# Login

ALLOWED_USERS = {
    "lucas.costa@mkthouse.com.br": "mudar12345",
    "gabriel.garcia@mkthouse.com.br": "Peter2025!",
    "daniela.scibor@mkthouse.com.br": "mudar12345",
    "regiane.paula@mkthouse.com.br": "mudar12345",
    "pamela.fructuoso@mkthouse.com.br": "mudar12345",
    "fernanda.sabino@mkthouse.com.br": "mudar12345",
    "cacia.nogueira@mkthouse.com.br": "mudar12345",
    "edson.fortaleza@mkthouse.com.br": "mudar12345",
    "lucas.depaula@mkthouse.com.br": "mudar12345",
    "janaina.morais@mkthouse.com.br": "mudar12345",
    "debora.ramos@mkthouse.com.br": "mudar12345",
    "david.silva@mkthouse.com.br": "mudar12345",
}
ALLOWED_USERS = {k.strip().lower(): v for k, v in ALLOWED_USERS.items()}

def do_login():
    st.markdown(LOGIN_CSS, unsafe_allow_html=True)
    st.title("🔒 Acesso Restrito")
    st.markdown('<div class="login-wrap">', unsafe_allow_html=True)
    st.markdown('<div class="login-hero">Use seu e-mail corporativo. Em caso de dúvidas, contate o BI.</div>', unsafe_allow_html=True)
    st.markdown('<div class="login-card">', unsafe_allow_html=True)
    with st.form("login_form", clear_on_submit=False):
        email = st.text_input("E-mail", placeholder="seu.email@mkthouse.com.br")
        pwd = st.text_input("Senha", type="password", placeholder="••••••••")
        entrar = st.form_submit_button("Entrar")
    st.markdown('</div></div>', unsafe_allow_html=True)

    if entrar:
        email_norm = (email or "").strip().lower()
        if email_norm in ALLOWED_USERS and pwd == ALLOWED_USERS[email_norm]:
            st.session_state.auth = True
            st.session_state.user_email = email_norm
            logger.info(f"Usuário autenticado: {email_norm}")
            st.rerun()
        else:
            logger.warning(f"Tentativa de login falhou: {email_norm}")
            st.error("Credenciais inválidas.")

if "auth" not in st.session_state:
    st.session_state.auth = False


# === PARTE 4/10 =====================================================
# Funções utilitárias (utils)

URL_RE = re.compile(r'https?://\S+')

def extrair_links(celula):
    if pd.isna(celula):
        return []
    t = str(celula).replace(",", " ").replace("(", " ").replace(")", " ").replace('"', " ").replace("'", " ")
    return [u.rstrip(").,") for u in URL_RE.findall(t)]

def redimensionar(img: Image.Image, max_w: int, max_h: int) -> Image.Image:
    img = ImageOps.exif_transpose(img)
    if img.mode != "RGB":
        img = img.convert("RGB")
    img.thumbnail((max_w, max_h), resample=Image.LANCZOS)
    return img

def comprimir_jpeg_binsearch(img: Image.Image, limite_kb: int) -> BytesIO:
    lo, hi, best = 35, 95, None
    buf = BytesIO()
    img.save(buf, "JPEG", quality=75, optimize=True, progressive=True, subsampling=2)
    if buf.tell()/1024 <= limite_kb:
        buf.seek(0)
        return buf
    best = buf
    while lo <= hi:
        mid = (lo+hi)//2
        buf = BytesIO()
        img.save(buf, "JPEG", quality=mid, optimize=True, progressive=True, subsampling=2)
        if buf.tell()/1024 <= limite_kb:
            best = buf
            lo = mid+1
        else:
            hi = mid-1
    if best is None:
        best = BytesIO()
        img.save(best, "JPEG", quality=35, optimize=True, progressive=True, subsampling=2)
    best.seek(0)
    return best

def px_to_inches(px): 
    return Inches(px / 96.0)

def hex_to_rgb(hex_str: str):
    s = hex_str.strip().lstrip("#")
    if len(s) == 3:
        s = "".join([c*2 for c in s])
    return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)

def pick_contrast_color(r, g, b):
    brightness = (r*299 + g*587 + b*114) / 1000
    return (0,0,0) if brightness > 128 else (255,255,255)

# --- HASHES p/ duplicatas ---
def _sha1_bytes(b: bytes) -> str:
    return hashlib.sha1(b).hexdigest()

def _img_dhash(img: Image.Image, hash_size: int = 8) -> str:
    """Perceptual hash (dHash) simples, estável e sem libs externas."""
    im = ImageOps.exif_transpose(img).convert("L").resize((hash_size + 1, hash_size), Image.LANCZOS)
    pixels = np.asarray(im, dtype=np.int16)
    diff = pixels[:, 1:] > pixels[:, :-1]
    bits = 0
    for row in diff:
        for v in row:
            bits = (bits << 1) | int(v)
    return f"{bits:0{hash_size*hash_size//4}x}"


# === PARTE 5/10 =====================================================
# Qualidade da imagem + Efeitos

def _laplacian_var_gray(pil_img: Image.Image) -> float:
    g = pil_img.convert("L")
    a = np.asarray(g, dtype=np.float32)
    H, W = a.shape
    if H < 3 or W < 3:
        return 0.0
    out = (a[0:-2,1:-1] + a[1:-1,0:-2] + a[1:-1,2:] + a[2:,1:-1] - 4*a[1:-1,1:-1])
    return float(out.var())

def medir_qualidade(img: Image.Image) -> dict:
    im = ImageOps.exif_transpose(img)
    w, h = im.size
    im_small = im.copy()
    im_small.thumbnail((1024, 1024), Image.LANCZOS)
    g = im_small.convert("L")
    arr = np.asarray(g, dtype=np.float32)
    return {
        "width": int(w),
        "height": int(h),
        "megapixels": float((w*h)/1_000_000.0),
        "mean_brightness": float(arr.mean()),
        "std_contrast": float(arr.std()),
        "blur_score": _laplacian_var_gray(im_small),
    }

def _hex_to_rgba_tuple(hex_color, alpha=255):
    s = hex_color.strip().lstrip("#")
    if len(s) == 3:
        s = "".join([c*2 for c in s])
    r, g, b = int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)
    return (r, g, b, alpha)

def _apply_rounded_corners(img_rgba: Image.Image, radius: int) -> Image.Image:
    if radius <= 0:
        return img_rgba
    w, h = img_rgba.size
    mask = Image.new("L", (w, h), 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle([0, 0, w, h], radius=radius, fill=255)
    out = img_rgba.copy()
    out.putalpha(mask)
    return out

def _apply_border_color(img_rgba: Image.Image, border_px: int, border_hex: str, radius: int) -> Image.Image:
    if border_px <= 0:
        return img_rgba
    w, h = img_rgba.size
    result = Image.new("RGBA", (w + 2*border_px, h + 2*border_px), (0,0,0,0))
    draw = ImageDraw.Draw(result)
    outer = [0, 0, result.size[0], result.size[1]]
    inner = [border_px, border_px, border_px + w, border_px + h]
    draw.rounded_rectangle(outer, radius=radius+border_px, fill=_hex_to_rgba_tuple(border_hex))
    hole = Image.new("L", result.size, 255)
    hole_draw = ImageDraw.Draw(hole)
    hole_draw.rounded_rectangle(inner, radius=radius, fill=0)
    result.putalpha(hole)
    result.alpha_composite(img_rgba, dest=(border_px, border_px))
    return result

def _apply_drop_shadow(img_rgba: Image.Image, blur: int, offset: int, opacity_pct: int) -> Image.Image:
    if blur <= 0 and offset <= 0:
        return img_rgba
    alpha = img_rgba.split()[-1]
    a = max(0, min(255, int(255 * (opacity_pct/100))))
    pad = blur + offset + 2
    w, h = img_rgba.size
    canvas = Image.new("RGBA", (w + pad, h + pad), (0,0,0,0))
    shadow = Image.new("RGBA", (w, h), (0,0,0,a))
    shadow.putalpha(alpha)
    shadow = shadow.filter(ImageFilter.GaussianBlur(radius=blur))
    canvas.alpha_composite(shadow, dest=(offset, offset))
    canvas.alpha_composite(img_rgba, dest=(0,0))
    return canvas

def apply_effects_pipeline(img_rgb: Image.Image, cfg: dict) -> Image.Image:
    out = img_rgb.convert("RGBA")
    if cfg.get("fx_round"):
        out = _apply_rounded_corners(out, int(cfg.get("fx_round_radius", 20)))
    if cfg.get("fx_border"):
        out = _apply_border_color(
            out,
            int(cfg.get("fx_border_width", 6)),
            cfg.get("fx_border_color", "#DDDDDD"),
            int(cfg.get("fx_round_radius", 20)) if cfg.get("fx_round") else 0
        )
    if cfg.get("fx_shadow"):
        out = _apply_drop_shadow(
            out,
            int(cfg.get("fx_shadow_blur", 10)),
            int(cfg.get("fx_shadow_offset", 8)),
            int(cfg.get("fx_shadow_opac", 40)),
        )
    return out


# === PARTE 6/10 =====================================================
# PPT helpers + Download & Process

def get_slots(n, prs):
    IMG_TOP = Inches(1.2); CONTENT_W = Inches(11); CONTENT_H = Inches(6); GAP = Inches(0.2)
    start_left = (prs.slide_width - CONTENT_W) / 2
    if n == 1:
        return [(start_left, IMG_TOP, CONTENT_W, CONTENT_H)]
    cols = n
    total_gap = GAP * (cols - 1)
    cell_w = (CONTENT_W - total_gap) / cols
    return [(start_left + c*(cell_w+GAP), IMG_TOP, cell_w, CONTENT_H) for c in range(cols)]

def add_title_and_address(slide, title_text, address_text, title_rgb=(0,0,0),
                          font_name="Radikal", title_font_size_pt=18, title_font_bold=True):
    TITLE_LEFT, TITLE_TOP, TITLE_W = Inches(0.5), Inches(0.2), Inches(12)
    tx = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_W, Inches(1))
    tf = tx.text_frame; tf.clear()
    p = tf.paragraphs[0]; run = p.add_run(); run.text = title_text
    f = run.font; f.name = font_name or "Radikal"; f.size = Pt(title_font_size_pt or 18)
    f.bold = bool(title_font_bold); f.color.rgb = RGBColor(*title_rgb)
    p.alignment = 1
    if address_text:
        p2 = tf.add_paragraph()
        run2 = p2.add_run(); run2.text = address_text
        f2 = run2.font; f2.name = font_name or "Radikal"; f2.size = Pt(max(8, (title_font_size_pt or 18) / 2))
        f2.bold = False; f2.color.rgb = RGBColor(*title_rgb)
        p2.alignment = 1

def set_slide_bg(slide, rgb_tuple):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = RGBColor(*rgb_tuple)

def place_picture(slide, buf, w_px, h_px, left, top, max_w_in, max_h_in):
    img_w_in = px_to_inches(w_px); img_h_in = px_to_inches(h_px)
    ratio = min(float(max_w_in)/float(img_w_in), float(max_h_in)/float(img_h_in), 1.0)
    final_w = img_w_in * ratio; final_h = img_h_in * ratio
    x = left + (max_w_in - final_w)/2; y = top + (max_h_in - final_h)/2
    buf.seek(0)
    slide.shapes.add_picture(buf, x, y, width=final_w, height=final_h)

def is_portrait(w_px: int, h_px: int, tol: float = 1.05) -> bool:
    if w_px <= 0 or h_px <= 0:
        return False
    return (h_px / float(w_px)) >= tol

def move_slide_to_index(prs, old_index, new_index):
    sldIdLst = prs.slides._sldIdLst
    sld = sldIdLst[old_index]
    sldIdLst.remove(sld)
    sldIdLst.insert(new_index, sld)

def add_logo_top_right(slide, prs, logo_bytes: bytes, logo_width_in: float):
    if not logo_bytes: 
        return
    left = prs.slide_width - Inches(0.5) - Inches(logo_width_in); top = Inches(0.2)
    slide.shapes.add_picture(BytesIO(logo_bytes), left, top, width=Inches(logo_width_in))

def add_signature_bottom_right(slide, prs, signature_bytes: bytes, signature_width_in: float,
                               bottom_margin_in: float = 0.2, right_margin_in: float = 0.2):
    if not signature_bytes: 
        return
    try:
        im = Image.open(BytesIO(signature_bytes)); w_px, h_px = im.size
        ratio = (h_px / float(w_px)) if w_px else 0.4
    except Exception:
        ratio = 0.4
    sig_h_in = signature_width_in * ratio
    left = prs.slide_width - Inches(right_margin_in) - Inches(signature_width_in)
    top  = prs.slide_height - Inches(bottom_margin_in) - Inches(sig_h_in)
    slide.shapes.add_picture(BytesIO(signature_bytes), left, top, width=Inches(signature_width_in))

# -------------------------------------------------------------------
# DOWNLOAD & PROCESS (com hashes SHA1 + dHash + logs)
# -------------------------------------------------------------------
def baixar_processar(session, url: str, max_w: int, max_h: int, limite_kb: int, timeout: int, fx_cfg: dict = None):
    try:
        logger.debug(f"Iniciando download: {url}")
        r = session.get(url, timeout=timeout, stream=True)
        if r.status_code != 200:
            logger.warning(f"HTTP {r.status_code} ao baixar: {url}")
            return (url, False, None, None, None)

        raw_bytes = r.content
        img = Image.open(BytesIO(raw_bytes))

        sha1_hex = _sha1_bytes(raw_bytes)
        im_small_for_hash = img.copy()
        im_small_for_hash.thumbnail((256, 256), Image.LANCZOS)
        dhash_hex = _img_dhash(im_small_for_hash)

        quality = medir_qualidade(img)  # métricas
        img = redimensionar(img, max_w, max_h)

        need_alpha = False
        if fx_cfg and (fx_cfg.get("fx_shadow") or fx_cfg.get("fx_round") or fx_cfg.get("fx_border")):
            img_rgba = apply_effects_pipeline(img.convert("RGB"), fx_cfg)
            need_alpha = True
        else:
            img_rgba = img.convert("RGBA")

        if need_alpha:
            buf = BytesIO(); img_rgba.save(buf, format="PNG", optimize=True)
            if buf.tell() / 1024 <= limite_kb:
                buf.seek(0); w, h = img_rgba.size
                return (url, True, buf, (w, h), quality, sha1_hex, dhash_hex)
            pal = img_rgba.convert("P", palette=Image.ADAPTIVE, colors=256)
            buf = BytesIO(); pal.save(buf, format="PNG", optimize=True)
            if buf.tell() / 1024 <= limite_kb:
                buf.seek(0); w, h = img_rgba.size
                return (url, True, buf, (w, h), quality, sha1_hex, dhash_hex)
            # fallback para JPEG com fundo branco
            bg = Image.new("RGB", img_rgba.size, (255, 255, 255))
            bg.paste(img_rgba, mask=img_rgba.split()[-1])
            buf = comprimir_jpeg_binsearch(bg, limite_kb)
            w, h = bg.size
            return (url, True, buf, (w, h), quality, sha1_hex, dhash_hex)
        else:
            buf = comprimir_jpeg_binsearch(img.convert("RGB"), limite_kb)
            w, h = img.size
            return (url, True, buf, (w, h), quality, sha1_hex, dhash_hex)

    except requests.exceptions.Timeout:
        logger.warning(f"Timeout ao baixar: {url}")
        return (url, False, None, None, None)
    except Exception as e:
        logger.warning(f"Falha ao processar {url}: {e}")
        return (url, False, None, None, None)


# === PARTE 7/10 =====================================================
# ZIP de imagens + PPT com modelo (capa/final)

def _buf_to_jpeg_bytes(img_buf: BytesIO) -> bytes:
    try:
        img_buf.seek(0)
        im = Image.open(img_buf)
        fmt = (im.format or "").upper()
        if fmt in ("JPEG", "JPG"):
            img_buf.seek(0)
            return img_buf.read()
        if im.mode in ("RGBA", "LA"):
            bg = Image.new("RGB", im.size, (255, 255, 255))
            bg.paste(im, mask=im.split()[-1])
            im = bg
        else:
            im = im.convert("RGB")
        out = BytesIO()
        im.save(out, "JPEG", quality=88, optimize=True, progressive=True, subsampling=2)
        out.seek(0)
        return out.read()
    except Exception as e:
        logger.warning(f"_buf_to_jpeg_bytes fallback: {e}")
        img_buf.seek(0)
        return img_buf.read()

def _sanitize_folder_name(name: str) -> str:
    safe = re.sub(r'[\\/:*?"<>|]+', ' ', str(name or "").strip())
    safe = re.sub(r'\s+', ' ', safe)
    return safe[:80] if len(safe) > 80 else safe

def montar_zip_imagens(items, resultados, excluded_urls: set) -> BytesIO:
    grupos = OrderedDict()
    for loja, endereco, url in items:
        if (url in resultados) and (url not in excluded_urls):
            grupos.setdefault(str(loja), []).append((url, resultados[url]))

    mem_zip = BytesIO()
    with zipfile.ZipFile(mem_zip, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
        for loja, lista in grupos.items():
            pasta = _sanitize_folder_name(loja) or "Sem Nome"
            contador = 1
            for _url, (_loja, _end, buf, (w, h), *_) in lista:
                jpeg_bytes = _buf_to_jpeg_bytes(buf)
                arquivo = f"{pasta} - {contador}.jpg"
                caminho = f"{pasta}/{arquivo}"
                try:
                    zf.writestr(caminho, jpeg_bytes)
                except Exception as e:
                    logger.warning(f"Falha ao escrever {caminho} no ZIP: {e}")
                contador += 1
    mem_zip.seek(0)
    return mem_zip

def gerar_ppt_modelo_capa_final(
    template_bytes: bytes,
    items, resultados, titulo,
    max_per_slide, sort_mode, bg_rgb,
    logo_bytes=None, logo_width_in=1.2,
    signature_bytes=None, signature_width_in=None, auto_half_signature=True,
    signature_bottom_margin_in=0.2, signature_right_margin_in=0.2,
    title_font_name="Radikal", title_font_size_pt=18, title_font_bold=True,
    excluded_urls=None
):
    excluded_urls = excluded_urls or set()
    prs = Presentation(BytesIO(template_bytes))
    logger.info(f"Template carregado: {len(prs.slides)} slides (esperado: 2 – capa e final)")

    final_idx = 1 if len(prs.slides) >= 2 else None
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    groups = OrderedDict()
    for loja, endereco, url in items:
        if url in resultados and url not in excluded_urls:
            groups.setdefault(str(loja), []).append((url, resultados[url]))

    if sort_mode == "Nome da loja (A→Z)":
        loja_keys = sorted(groups.keys(), key=lambda s: (s is None or str(s).strip()== "", (s or "").strip().casefold()))
    else:
        loja_keys = list(groups.keys())

    title_rgb = pick_contrast_color(*bg_rgb)
    signature_width = (logo_width_in/2.0) if auto_half_signature else (signature_width_in or 0.6)

    for loja in loja_keys:
        imgs = groups[loja]
        i = 0
        while i < len(imgs):
            if max_per_slide == "Automático":
                _url0, (_loja0, endereco, _buf0, (w0, h0), *_r0) = imgs[i]
                per_slide = 3 if is_portrait(w0, h0) else 2
            else:
                endereco = imgs[i][1][1]
                per_slide = int(max_per_slide)

            batch = imgs[i:i+per_slide]
            i += per_slide
            slide = prs.slides.add_slide(blank_layout)
            set_slide_bg(slide, bg_rgb)
            add_title_and_address(slide, loja, endereco, title_rgb,
                                  title_font_name, title_font_size_pt, title_font_bold)
            if logo_bytes:
                add_logo_top_right(slide, prs, logo_bytes, logo_width_in or 1.2)
            if signature_bytes:
                add_signature_bottom_right(
                    slide, prs, signature_bytes, signature_width,
                    bottom_margin_in=signature_bottom_margin_in,
                    right_margin_in=signature_right_margin_in
                )
            slots = get_slots(len(batch), prs)
            for (url, (_loja, _end, buf, (w_px, h_px), *rest)), (left, top, max_w_in, max_h_in) in zip(batch, slots):
                try:
                    place_picture(slide, buf, w_px, h_px, left, top, max_w_in, max_h_in)
                except Exception as e:
                    logger.warning(f"Falha ao inserir imagem no slide ({url}): {e}")

    if final_idx is not None and final_idx < len(prs.slides):
        try:
            move_slide_to_index(prs, final_idx, len(prs.slides)-1)
        except Exception as e:
            logger.warning(f"Não foi possível mover o slide final: {e}")

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    logger.info("PPT com modelo gerado com sucesso.")
    return out


# === PARTE 8/10 =====================================================
# Prévia de slides — composição PIL + batching

def _get_slots_px(n, canvas_w=1280, canvas_h=720):
    IMG_TOP = int(1.2 * 96)
    CONTENT_W = int(11 * 96)
    CONTENT_H = int(6 * 96)
    GAP = int(0.2 * 96)
    start_left = (canvas_w - CONTENT_W) // 2
    if n == 1:
        return [(start_left, IMG_TOP, CONTENT_W, CONTENT_H)]
    cols = n
    total_gap = GAP * (cols - 1)
    cell_w = (CONTENT_W - total_gap) // cols
    return [(start_left + c*(cell_w+GAP), IMG_TOP, cell_w, CONTENT_H) for c in range(cols)]

def _paste_fit(canvas: Image.Image, im: Image.Image, slot):
    left, top, max_w, max_h = slot
    w, h = im.size
    if w <= 0 or h <= 0:
        logger.warning("Imagem com dimensão inválida para prévia.")
        return
    ratio = min(max_w/float(w), max_h/float(h), 1.0)
    new_w = max(1, int(w*ratio))
    new_h = max(1, int(h*ratio))
    im2 = im.resize((new_w, new_h), Image.LANCZOS)
    x = left + (max_w - new_w)//2
    y = top + (max_h - new_h)//2
    canvas.paste(im2, (x, y))

def compose_slide_preview(batch, loja, endereco, cfg, canvas_w=1280, canvas_h=720):
    bg = cfg["bg_rgb"]
    title_rgb = pick_contrast_color(*bg)
    canvas = Image.new("RGB", (canvas_w, canvas_h), bg)

    try:
        font = ImageFont.truetype("arial.ttf", 28)
        font2 = ImageFont.truetype("arial.ttf", 16)
    except Exception:
        font = ImageFont.load_default()
        font2 = ImageFont.load_default()

    draw = ImageDraw.Draw(canvas)
    title = str(loja)

    try:
        _, _, w_title, h_title = draw.textbbox((0, 0), title, font=font)
    except Exception:
        w_title, h_title = 0, 0
    draw.text(((canvas_w - w_title)//2, 16), title, fill=title_rgb, font=font)

    if endereco:
        try:
            _, _, w_addr, h_addr = draw.textbbox((0, 0), endereco, font=font2)
        except Exception:
            w_addr, h_addr = 0, 0
        draw.text(((canvas_w - w_addr)//2, 16 + h_title + 6), endereco, fill=title_rgb, font=font2)

    slots = _get_slots_px(len(batch), canvas_w, canvas_h)
    for (url, (_loja, _end, buf, (w_px, h_px), *rest)), slot in zip(batch, slots):
        try:
            buf.seek(0)
            im = Image.open(buf).convert("RGB")
            _paste_fit(canvas, im, slot)
        except Exception as e:
            logger.warning(f"Falha ao compor imagem na prévia ({url}): {e}")
            l, t, sw, sh = slot
            draw.rectangle([l, t, l+sw, t+sh], outline=(200,0,0), width=3)
            draw.line([l, t, l+sw, t+sh], fill=(200,0,0), width=3)
            draw.line([l+sw, t, l, t+sh], fill=(200,0,0), width=3)

    return canvas

def build_batches_for_preview(items, resultados, cfg, excluded_urls):
    groups = OrderedDict()
    for loja, endereco, url in items:
        if url in resultados and url not in excluded_urls:
            groups.setdefault(str(loja), []).append((url, resultados[url]))

    if cfg["sort_mode"] == "Nome da loja (A→Z)":
        loja_keys = sorted(
            groups.keys(),
            key=lambda s: (s is None or str(s).strip()== "", (s or "").strip().casefold())
        )
    else:
        loja_keys = list(groups.keys())

    batches = []
    for loja in loja_keys:
        imgs = groups[loja]
        i = 0
        while i < len(imgs):
            if cfg["max_per_slide"] == "Automático":
                _url0, (_loja0, end0, _buf0, (w0, h0), *_r0) = imgs[i]
                per_slide = 3 if is_portrait(w0, h0) else 2
            else:
                end0 = imgs[i][1][1]
                per_slide = int(cfg["max_per_slide"])
            batch = imgs[i:i+per_slide]
            i += per_slide
            batches.append((loja, end0, batch))
    return batches


# === PARTE 9/10 =====================================================
# UI de miniaturas + detecção + reset

def img_to_html_with_border(image: Image.Image, width_px: int, border_px: int, border_color: str):
    im = image.copy()
    im.thumbnail((width_px, width_px))
    buf = BytesIO()
    im.save(buf, format="PNG")
    b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
    style = (
        f"border:{border_px}px solid {border_color};"
        f"border-radius:10px;display:block;max-width:100%;width:{width_px}px;"
    )
    return f'<img src="data:image/png;base64,{b64}" style="{style}" />'

def render_steps(current: int):
    labels = ["Upload", "Pré-visualização", "Gerar/Exportar"]
    html = ['<div class="steps">']
    for i, txt in enumerate(labels, start=1):
        cls = "step active" if i == current else "step"
        html.append(f'<div class="{cls}">{i}. {txt}</div>')
        if i < len(labels):
            html.append('<span class="step sep">›</span>')
    html.append("</div>")
    st.markdown("".join(html), unsafe_allow_html=True)

def render_summary(items, resultados, excluded):
    total_urls = len(items)
    baixadas = sum(1 for _, _, url in items if url in resultados)
    lojas = len({loja for loja, _, _ in items})
    st.markdown(
        f"**Resumo:** "
        f"<span class='badge'>Lojas: {lojas}</span> "
        f"<span class='badge'>Links: {total_urls}</span> "
        f"<span class='badge'>Baixadas: {baixadas}</span> "
        f"<span class='badge'>Excluídas: {len(excluded)}</span>",
        unsafe_allow_html=True
    )

def _cb_select_all(loja, imgs):
    excluded = st.session_state.excluded_urls
    for url, _ in imgs:
        excluded.add(url)
    st.rerun()

def _cb_clear_group(loja, imgs):
    excluded = st.session_state.excluded_urls
    for url, _ in imgs:
        excluded.discard(url)
    st.rerun()

def render_preview(items, resultados, sort_mode, thumb_px: int, thumbs_per_row: int):
    if "excluded_urls" not in st.session_state:
        st.session_state.excluded_urls = set()
    if "expanded_groups" not in st.session_state:
        st.session_state.expanded_groups = {}

    excluded = st.session_state.excluded_urls
    expanded_groups = st.session_state.expanded_groups

    groups = OrderedDict()
    for loja, endereco, url in items:
        if url in resultados:
            groups.setdefault(str(loja), []).append((url, resultados[url]))

    if sort_mode == "Nome da loja (A→Z)":
        loja_keys = sorted(
            groups.keys(),
            key=lambda s: (s is None or str(s).strip()== "", (s or "").strip().casefold())
        )
    else:
        loja_keys = list(groups.keys())

    low_q = set(st.session_state.get("low_quality_urls", []))
    dups = set(st.session_state.get("duplicate_urls", []))

    c0a, c0b = st.columns([1,1])
    with c0a:
        if st.button(
            f"⚠️ Excluir todas de BAIXA QUALIDADE ({len(low_q)})",
            type="secondary", use_container_width=True, disabled=(len(low_q)==0)
        ):
            st.session_state.excluded_urls.update(low_q)
            st.rerun()
    with c0b:
        if st.button(
            f"🔁 Excluir todas DUPLICADAS ({len(dups)})",
            type="secondary", use_container_width=True, disabled=(len(dups)==0)
        ):
            st.session_state.excluded_urls.update(dups)
            st.rerun()

    c1, c2, c3, c4 = st.columns([1,1,1,1])
    with c1:
        if st.button("🧹 Limpar todas as exclusões", type="secondary", use_container_width=True):
            excluded.clear()
            st.rerun()
    with c2:
        if st.button("🔁 Inverter seleção", type="secondary", use_container_width=True):
            all_urls = {url for _, v in groups.items() for (url, _) in v}
            st.session_state.excluded_urls = all_urls - excluded
            st.rerun()
    with c3:
        if st.button("➕ Expandir todas", type="secondary", use_container_width=True):
            for loja in loja_keys:
                expanded_groups[loja] = True
            st.rerun()
    with c4:
        if st.button("➖ Recolher todas", type="secondary", use_container_width=True):
            for loja in loja_keys:
                expanded_groups[loja] = False
            st.rerun()

    st.caption(f"Excluídas até agora: **{len(excluded)}**")

    for loja in loja_keys:
        imgs = groups[loja]
        key_toggle = f"tg_{hash(loja)}"
        current_state = expanded_groups.get(loja, True)

        st.markdown(
            f"""<div class="group-head">
                <div><strong>📄 {loja}</strong> — {len(imgs)} foto(s)</div>
                <div class="badge">{'Aberto' if current_state else 'Fechado'}</div>
            </div>""",
            unsafe_allow_html=True
        )
        new_state = st.toggle("Manter este grupo visível", value=current_state, key=key_toggle)
        expanded_groups[loja] = bool(new_state)

        gc1, gc2 = st.columns([1,1])
        with gc1:
            st.button(
                f"Selecionar todas de {loja}",
                key=f"sel_all_{hash(loja)}",
                use_container_width=True,
                on_click=_cb_select_all,
                args=(loja, imgs)
            )
        with gc2:
            st.button(
                f"Limpar seleção de {loja}",
                key=f"clr_sel_{hash(loja)}",
                use_container_width=True,
                on_click=_cb_clear_group,
                args=(loja, imgs)
            )

        if expanded_groups[loja]:
            cols = st.columns(thumbs_per_row)
            col_idx = 0
            for (url, (_loja, _end, buf, (w_px, h_px), *rest)) in imgs:
                with cols[col_idx]:
                    try:
                        buf.seek(0)
                        im = Image.open(buf)
                    except Exception:
                        st.warning("Não foi possível pré-visualizar esta imagem.")
                        col_idx = (col_idx + 1) % thumbs_per_row
                        continue

                    is_excluded = url in excluded
                    border_px = 3 if is_excluded else 1
                    border_color = "#E53935" if is_excluded else "#DDDDDD"
                    st.markdown('<div class="img-card">', unsafe_allow_html=True)
                    st.markdown(
                        img_to_html_with_border(im, thumb_px, border_px, border_color),
                        unsafe_allow_html=True
                    )

                    badges = []
                    if url in low_q:
                        badges.append("<span class='badge' style='background:#FDECEA;color:#B71C1C'>Baixa qualidade</span>")
                    if url in dups:
                        badges.append("<span class='badge' style='background:#FFF8E1;color:#E65100'>Possível duplicata</span>")
                    if badges:
                        st.markdown(" ".join(badges), unsafe_allow_html=True)

                    key = "ex_" + hashlib.md5(url.encode("utf-8")).hexdigest()
                    checked = st.checkbox("Excluir esta foto", key=key, value=is_excluded)
                    if checked:
                        excluded.add(url)
                    else:
                        excluded.discard(url)
                    st.markdown('</div>', unsafe_allow_html=True)
                col_idx = (col_idx + 1) % thumbs_per_row
        st.divider()

def detectar_problemas(resultados, min_mp=0.8, min_blur=45):
    low_quality = set()
    by_sha1, by_dhash = {}, {}
    for url, tup in resultados.items():
        loja, endereco, buf, (w, h), quality, sha1_hex, dhash_hex = tup
        if (quality.get("megapixels", 0) < min_mp) or (quality.get("blur_score", 0) < min_blur):
            low_quality.add(url)
        by_sha1.setdefault(sha1_hex, []).append(url)
        by_dhash.setdefault(dhash_hex, []).append(url)

    duplicates = set()
    for group in by_sha1.values():
        if len(group) > 1:
            duplicates.update(group[1:])
    for group in by_dhash.values():
        if len(group) > 1:
            duplicates.update([u for u in group[1:] if u not in duplicates])

    logger.info(f"detectar_problemas: low={len(low_quality)} dup={len(duplicates)}")
    return low_quality, duplicates

def reset_app(preserve_login: bool = True):
    user = st.session_state.get("user_email")
    auth = st.session_state.get("auth", False)

    st.session_state.clear()

    st.session_state.xlsx_key = 0
    st.session_state.template_key = 0
    st.session_state.logo_key = 0
    st.session_state.sign_key = 0
    st.session_state.download_key = 0
    st.session_state.images_zip_key = 0

    st.session_state.exp_plan = True
    st.session_state.exp_style = False
    st.session_state.exp_brand = False
    st.session_state.exp_fx = False
    st.session_state.exp_perf = False
    st.session_state.exp_model = False

    st.session_state.ppt_bytes = None
    st.session_state.images_zip_bytes = None
    st.session_state.generated = False
    st.session_state.output_filename = "Modelo_01"

    if preserve_login and auth:
        st.session_state.auth = True
        st.session_state.user_email = user
        st.session_state.dark_mode = False
    st.rerun()


# === PARTE 10/10 =====================================================
# APP (main_app) + Roteamento final

def main_app():
    # Inicializações seguras
    for k in ["xlsx_key", "template_key", "logo_key", "sign_key", "download_key", "images_zip_key"]:
        if k not in st.session_state:
            st.session_state[k] = 0
    if "exp_plan" not in st.session_state:
        st.session_state.exp_plan = True
        st.session_state.exp_style = False
        st.session_state.exp_brand = False
        st.session_state.exp_fx = False
        st.session_state.exp_perf = False
        st.session_state.exp_model = False
    if "pipeline" not in st.session_state: st.session_state.pipeline = {}
    if "excluded_urls" not in st.session_state: st.session_state.excluded_urls = set()
    if "preview_mode" not in st.session_state: st.session_state.preview_mode = False
    if "expanded_groups" not in st.session_state: st.session_state.expanded_groups = {}
    if "output_filename" not in st.session_state: st.session_state.output_filename = "Modelo_01"
    if "generated" not in st.session_state: st.session_state.generated = False
    if "ppt_bytes" not in st.session_state: st.session_state.ppt_bytes = None
    if "images_zip_bytes" not in st.session_state: st.session_state.images_zip_bytes = None
    if "preview_bump" not in st.session_state: st.session_state.preview_bump = 0

    with st.sidebar:
        st.header("⚙️ Preferências")
        st.session_state.dark_mode = st.toggle("Usar tema escuro", value=st.session_state.dark_mode)
        apply_theme(st.session_state.dark_mode)

        with st.expander("📄 Planilha & Layout", expanded=st.session_state.exp_plan):
            st.caption("Colunas (nomes exatos do cabeçalho):")
            loja_col = st.text_input("🛒 Coluna de LOJA", value="Selecione sua loja", key="loja_col")
            img_col  = st.text_input("🖼️ Coluna de FOTOS", value="Faça o upload das fotos", key="img_col")
            use_address = st.checkbox("➕ Incluir endereço abaixo do nome da loja", value=False, key="use_address")
            address_col = st.text_input("🏠 Coluna de ENDEREÇO", value="Endereço", key="address_col", disabled=not use_address)
            max_per_slide = st.selectbox("📐 Fotos por slide (máx.)", ["Automático", 1, 2, 3], index=0, key="max_per_slide")
            sort_mode = st.selectbox("🔤 Ordenar lojas por", ["Ordem original do Excel", "Nome da loja (A→Z)"], index=0, key="sort_mode")

        with st.expander("🎨 Aparência do slide", expanded=st.session_state.exp_style):
            bg_hex = st.color_picker("🎨 Cor de fundo", value="#FFFFFF", key="bg_hex")
            st.caption("Título do slide")
            title_font_name = st.text_input("Fonte do título", value="Radikal", key="title_font_name")
            title_font_size_pt = st.slider("Tamanho (pt)", 8, 48, 18, 1, key="title_font_size_pt")
            title_font_bold = st.checkbox("Negrito", value=True, key="title_font_bold")

        with st.expander("🏷️ Logo & ✍️ Assinatura", expanded=st.session_state.exp_brand):
            st.caption("Logo (canto superior direito)")
            logo_file = st.file_uploader("Logo (PNG/JPG)", type=["png", "jpg", "jpeg"], key=f"logo_uploader_{st.session_state.logo_key}")
            if "logo_bytes" not in st.session_state: st.session_state.logo_bytes = None
            if logo_file is not None: st.session_state.logo_bytes = logo_file.getvalue()
            logo_width_in = st.slider("Largura do LOGO (pol)", 0.5, 3.0, 1.2, 0.1, key="logo_width_in")

            st.markdown("---")
            st.caption("Assinatura (canto inferior direito)")
            signature_file = st.file_uploader("Assinatura (PNG/JPG)", type=["png", "jpg", "jpeg"], key=f"signature_uploader_{st.session_state.sign_key}")
            if "signature_bytes" not in st.session_state: st.session_state.signature_bytes = None
            if signature_file is not None: st.session_state.signature_bytes = signature_file.getvalue()

            auto_half_signature = st.checkbox("Usar 1/2 do tamanho do logo (recomendado)", value=True, key="auto_half_signature")
            derived_default_sig = (st.session_state.get("logo_width_in", 1.2) / 2.0)
            if not auto_half_signature:
                signature_width_in = st.slider("Largura da assinatura (pol)", 0.3, 2.0, float(derived_default_sig), 0.05, key="signature_width_in")
            else:
                if "signature_width_in" not in st.session_state:
                    st.session_state.signature_width_in = float(derived_default_sig)

            st.caption("Posição da assinatura (quanto menor, mais encostada no canto)")
            signature_right_margin_in = st.slider("Margem direita (pol)", 0.0, 1.0, 0.20, 0.05, key="sig_right_margin")
            signature_bottom_margin_in = st.slider("Margem inferior (pol)", 0.0, 1.0, 0.20, 0.05, key="sig_bottom_margin")

        with st.expander("📑 Modelo (capa + final)", expanded=st.session_state.exp_model):
            use_template = st.checkbox("Usar modelo (Capa + Final)", value=False, key="use_template")
            template_file = None
            if use_template:
                template_file = st.file_uploader("Suba o PPTX com 2 slides (Capa e Final)", type=["pptx"], key=f"template_pptx_{st.session_state.template_key}")

        with st.expander("🎭 Efeitos nas fotos", expanded=st.session_state.exp_fx):
            st.caption("Ative efeitos opcionais.")
            fx_shadow = st.checkbox("Sombra projetada", value=False, key="fx_shadow")
            shadow_blur = st.slider("Intensidade da sombra (blur)", 0, 30, 10, 1, key="fx_shadow_blur", disabled=not fx_shadow)
            shadow_offset = st.slider("Deslocamento da sombra (px)", 0, 30, 8, 1, key="fx_shadow_offset", disabled=not fx_shadow)
            shadow_opacity = st.slider("Opacidade da sombra (%)", 10, 100, 40, 5, key="fx_shadow_opac", disabled=not fx_shadow)
            fx_round = st.checkbox("Borda arredondada", value=False, key="fx_round")
            round_radius = st.slider("Raio dos cantos (px)", 0, 60, 20, 2, key="fx_round_radius", disabled=not fx_round)
            fx_border = st.checkbox("Borda colorida", value=False, key="fx_border")
            border_color_hex = st.color_picker("Cor da borda", value="#DDDDDD", key="fx_border_color", disabled=not fx_border)
            border_width = st.slider("Espessura da borda (px)", 1, 30, 6, 1, key="fx_border_width", disabled=not fx_border)

        with st.expander("⚡ Performance & Qualidade", expanded=st.session_state.exp_perf):
            thumb_px = st.slider("Tamanho das miniaturas (px)", 120, 400, 220, 10, key="thumb_px")
            thumbs_per_row = st.slider("Miniaturas por linha", 2, 8, 4, 1, key="thumbs_per_row")
            st.caption("Redimensionamento / compressão")
            target_w = st.number_input("Largura máx (px)", 480, 4096, 1280, 10, key="target_w")
            target_h = st.number_input("Altura máx (px)",  360, 4096, 720, 10, key="target_h")
            limite_kb = st.number_input("Tamanho máx por foto (KB)", 50, 2000, 450, 10, key="limite_kb")
            st.caption("Rede e paralelismo")
            max_workers = st.slider("Trabalhos em paralelo", 2, 32, 12, key="max_workers")
            req_timeout = st.slider("Timeout por download (s)", 5, 60, 15, key="req_timeout")
            st.caption("Critérios de qualidade (após o download)")
            min_mp = st.slider("Megapixels mínimos", 0.1, 5.0, 0.8, 0.1, key="min_megapixels")
            min_blur = st.slider("Limiar de nitidez (blur score)", 5, 300, 45, 5, key="min_blur_score")

    # Topo
    top_l, top_m, top_r = st.columns([5,1,1])
    with top_l:
        current_step = 1
        if st.session_state.get("preview_mode") and not st.session_state.get("generated"):
            current_step = 2
        if st.session_state.get("generated") or st.session_state.get("images_zip_bytes"):
            current_step = 3
        st.title("📸 Gerador de Book")
        render_steps(current_step)
        st.caption("Monte um PPT com fotos por loja (capa e final opcionais via modelo).")
    with top_m:
        st.markdown('<div class="reset-zone">', unsafe_allow_html=True)
        if st.button("Resetar", key="reset_btn", use_container_width=True, type="secondary"):
            st.session_state.xlsx_key += 1
            st.session_state.template_key += 1
            st.session_state.logo_key += 1
            st.session_state.sign_key += 1
            st.session_state.download_key += 1
            st.session_state.images_zip_key += 1
            reset_app(preserve_login=True)
        st.markdown('</div>', unsafe_allow_html=True)
    with top_r:
        st.markdown('<div class="logout-zone">', unsafe_allow_html=True)
        if st.button("Sair", key="logout_btn", use_container_width=True, type="secondary"):
            reset_app(preserve_login=False)
        st.markdown('</div>', unsafe_allow_html=True)

    # 1) Upload
    with st.expander("1. Upload", expanded=not st.session_state.preview_mode):
        up = st.file_uploader("Selecione a planilha (.xlsx)", type=["xlsx"], key=f"xlsx_upload_{st.session_state.xlsx_key}")
        btn_preview = st.button("🔎 Pré-visualizar", key="btn_preview", use_container_width=True)

        if btn_preview:
            if not up:
                st.warning("Envie a planilha primeiro.")
            else:
                try:
                    df = pd.read_excel(up)
                except Exception as e:
                    st.error(f"Não consegui ler o Excel: {e}")
                    st.stop()

                loja_col = st.session_state["loja_col"]
                img_col  = st.session_state["img_col"]
                use_address = st.session_state.get("use_address", False)
                address_col = st.session_state.get("address_col", "Endereço")

                required_cols = [loja_col, img_col] + ([address_col] if use_address else [])
                missing = [c for c in required_cols if c not in df.columns]
                if missing:
                    st.error(f"Colunas não encontradas: {missing}")
                    st.stop()

                items = []
                for _, row in df.iterrows():
                    loja = str(row[loja_col]).strip()
                    endereco = str(row[address_col]).strip() if use_address else ""
                    for url in extrair_links(row.get(img_col, "")):
                        if url.startswith("http"):
                            items.append((loja, endereco, url))

                seen, uniq = set(), []
                for loja, endereco, url in items:
                    if url not in seen:
                        seen.add(url)
                        uniq.append((loja, endereco, url))
                items = uniq

                if st.session_state["sort_mode"] == "Nome da loja (A→Z)":
                    grouped_tmp = OrderedDict()
                    for loja, end, url in items:
                        grouped_tmp.setdefault(loja, []).append((end, url))
                    items = [
                        (loja, end, url)
                        for loja in sorted(grouped_tmp.keys(),
                            key=lambda s: (s is None or str(s).strip()== "", (s or "").strip().casefold()))
                        for (end, url) in grouped_tmp[loja]
                    ]

                total = len(items)
                if total == 0:
                    st.warning("Nenhuma URL de imagem encontrada.")
                    st.stop()

                st.info(f"Baixando e processando **{total}** imagem(ns)...")
                session = requests.Session()
                adapter = requests.adapters.HTTPAdapter(
                    pool_connections=st.session_state["max_workers"],
                    pool_maxsize=st.session_state["max_workers"],
                    max_retries=2
                )
                session.mount("http://", adapter)
                session.mount("https://", adapter)
                session.headers.update({"User-Agent": "Mozilla/5.0 (GeradorBook Streamlit)"})

                fx_cfg = {
                    "fx_shadow": st.session_state.get("fx_shadow", False),
                    "fx_shadow_blur": st.session_state.get("fx_shadow_blur", 10),
                    "fx_shadow_offset": st.session_state.get("fx_shadow_offset", 8),
                    "fx_shadow_opac": st.session_state.get("fx_shadow_opac", 40),
                    "fx_round": st.session_state.get("fx_round", False),
                    "fx_round_radius": st.session_state.get("fx_round_radius", 20),
                    "fx_border": st.session_state.get("fx_border", False),
                    "fx_border_color": st.session_state.get("fx_border_color", "#DDDDDD"),
                    "fx_border_width": st.session_state.get("fx_border_width", 6),
                }

                prog = st.progress(0)
                status = st.empty()
                resultados, falhas, done = {}, 0, 0
                with ThreadPoolExecutor(max_workers=st.session_state["max_workers"]) as ex:
                    futures = {
                        ex.submit(
                            baixar_processar, session, url,
                            st.session_state["target_w"], st.session_state["target_h"],
                            st.session_state["limite_kb"], st.session_state["req_timeout"],
                            fx_cfg
                        ): (loja, endereco, url)
                        for loja, endereco, url in items
                    }
                    for fut in as_completed(futures):
                        loja, endereco, url = futures[fut]
                        try:
                            res = fut.result()
                        except Exception as e:
                            logger.error(f"Erro em download {url}: {e}")
                            res = (url, False, None, None, None)
                        if res[1] is True:
                            _url, _ok, buf, wh, quality, sha1_hex, dhash_hex = res
                            resultados[url] = (loja, endereco, buf, wh, quality, sha1_hex, dhash_hex)
                        else:
                            falhas += 1
                        done += 1
                        prog.progress(int(done * 100 / total))
                        status.write(f"Processadas {done}/{total} imagens...")

                status.write(f"Concluído. Falhas: {falhas}")

                low_q, dups = detectar_problemas(
                    resultados,
                    st.session_state["min_megapixels"],
                    st.session_state["min_blur_score"]
                )
                st.session_state.low_quality_urls = list(low_q)
                st.session_state.duplicate_urls = list(dups)

                st.session_state.pipeline = {
                    "items": items, "resultados": resultados, "falhas": falhas,
                    "settings": {
                        "max_per_slide": st.session_state["max_per_slide"],
                        "sort_mode": st.session_state["sort_mode"],
                        "bg_rgb": hex_to_rgb(st.session_state["bg_hex"]),
                        "title_font_name": st.session_state["title_font_name"],
                        "title_font_size_pt": st.session_state["title_font_size_pt"],
                        "title_font_bold": st.session_state["title_font_bold"],
                        "logo_bytes": st.session_state.logo_bytes,
                        "logo_width_in": st.session_state["logo_width_in"],
                        "signature_bytes": st.session_state.signature_bytes,
                        "auto_half_signature": st.session_state.get("auto_half_signature", True),
                        "signature_width_in": st.session_state.get("signature_width_in", (st.session_state["logo_width_in"]/2.0)),
                        "signature_right_margin_in": st.session_state.get("sig_right_margin", 0.20),
                        "signature_bottom_margin_in": st.session_state.get("sig_bottom_margin", 0.20),
                        "thumb_px": st.session_state["thumb_px"],
                        "thumbs_per_row": st.session_state["thumbs_per_row"],
                        "effects": fx_cfg,
                        "use_template": st.session_state.get("use_template", False),
                        "template_bytes": template_file.getvalue()
                            if (st.session_state.get("use_template") and template_file) else None,
                        "low_quality_urls": list(low_q),
                        "duplicate_urls": list(dups),
                    }
                }
                st.session_state.preview_mode = True
                st.session_state.generated = False
                st.session_state.ppt_bytes = None
                st.session_state.images_zip_bytes = None

                st.session_state.exp_plan = False
                st.session_state.exp_style = False
                st.session_state.exp_brand = False
                st.session_state.exp_fx = False
                st.session_state.exp_perf = False
                st.session_state.exp_model = False
                st.rerun()

    # 2) Pré-visualização
    with st.expander("2. Pré-visualização", expanded=st.session_state.preview_mode):
        if st.session_state.preview_mode and st.session_state.pipeline:
            p = st.session_state.pipeline
            render_summary(p["items"], p["resultados"], st.session_state.excluded_urls)
            render_preview(
                p["items"], p["resultados"],
                p["settings"]["sort_mode"],
                p["settings"]["thumb_px"],
                p["settings"]["thumbs_per_row"]
            )
            st.info("Marque **Excluir esta foto** nas imagens que não devem ir para o PPT/ZIP. Depois, avance para a etapa 3.")

    # 3) Gerar / Exportar  (👈 agora DENTRO da função e protegido com get)
    with st.expander("3. Gerar / Exportar", expanded=st.session_state.get("preview_mode", False)):
        if not (st.session_state.preview_mode and st.session_state.pipeline):
            st.info("Faça a pré-visualização antes de gerar/exportar.")
        else:
            cfg = st.session_state.pipeline["settings"]
            items = st.session_state.pipeline["items"]
            resultados = st.session_state.pipeline["resultados"]

            r1, r2 = st.columns([1, 5])
            with r1:
                if st.button("🔄 Recriar Prévia", use_container_width=True):
                    st.session_state.preview_bump = st.session_state.get("preview_bump", 0) + 1
                    st.rerun()
            with r2:
                st.caption("Gera novamente as miniaturas dos slides com base nas exclusões/ordem atuais.")

            batches = build_batches_for_preview(items, resultados, cfg, st.session_state.excluded_urls)
            st.subheader("🔎 Prévia dos primeiros slides")
            if len(batches) == 0:
                st.warning("Nenhum slide seria gerado com as configurações atuais.")
            else:
                preview_count = min(3, len(batches))
                cols = st.columns(preview_count)
                for idx in range(preview_count):
                    loja, end, batch = batches[idx]
                    canvas = compose_slide_preview(batch, loja, end, cfg).convert("RGBA")
                    W, H = canvas.width, canvas.height
                    title_rgb = pick_contrast_color(*cfg["bg_rgb"])

                    draw = ImageDraw.Draw(canvas)
                    top_bar_h = int(110)
                    draw.rectangle([0, 0, W, top_bar_h], fill=cfg["bg_rgb"])
                    try:
                        font_title = ImageFont.truetype("arial.ttf", 28)
                        font_addr  = ImageFont.truetype("arial.ttf", 16)
                    except Exception:
                        font_title = ImageFont.load_default()
                        font_addr  = ImageFont.load_default()

                    slide_w_in, slide_h_in = 13.33, 7.5
                    x_left = int(W * (0.30 / slide_w_in))
                    y_top  = int(H * (0.20 / slide_h_in))
                    draw.text((x_left, y_top), str(loja), fill=title_rgb, font=font_title, anchor="la")
                    if end:
                        try:
                            _, _, tw, th = draw.textbbox((0, 0), str(loja), font=font_title)
                            y_addr = y_top + th + 6
                        except Exception:
                            y_addr = y_top + 32
                        draw.text((x_left, y_addr), str(end), fill=title_rgb, font=font_addr, anchor="la")

                    if cfg.get("logo_bytes"):
                        try:
                            _bio = BytesIO(cfg["logo_bytes"])
                            logo_im = Image.open(_bio).convert("RGBA")
                            target_w_px = max(1, int(W * (cfg["logo_width_in"] / slide_w_in)))
                            ratio = target_w_px / float(logo_im.width)
                            logo_im = logo_im.resize((target_w_px, int(logo_im.height * ratio)), Image.LANCZOS)
                            x = W - int(W * (0.50 / slide_w_in)) - logo_im.width
                            y = int(H * (0.20 / slide_h_in))
                            canvas.alpha_composite(logo_im, (max(0, x), max(0, y)))
                        except Exception:
                            pass

                    if cfg.get("signature_bytes"):
                        try:
                            _bio = BytesIO(cfg["signature_bytes"])
                            sig_im = Image.open(_bio).convert("RGBA")
                            if cfg.get("auto_half_signature", True):
                                sig_w_in = (cfg.get("logo_width_in", 1.2) / 2.0)
                            else:
                                sig_w_in = cfg.get("signature_width_in") or 0.6
                            target_w_px = max(1, int(W * (sig_w_in / slide_w_in)))
                            ratio = target_w_px / float(sig_im.width)
                            sig_im = sig_im.resize((target_w_px, int(sig_im.height * ratio)), Image.LANCZOS)
                            right_margin_in  = float(cfg.get("signature_right_margin_in", 0.20))
                            bottom_margin_in = float(cfg.get("signature_bottom_margin_in", 0.20))
                            x = W - int(W * (right_margin_in / slide_w_in)) - sig_im.width
                            y = H - int(H * (bottom_margin_in / slide_h_in)) - sig_im.height
                            canvas.alpha_composite(sig_im, (max(0, x), max(0, y)))
                        except Exception:
                            pass

                    with cols[idx]:
                        st.image(canvas.convert("RGB"), caption=f"Slide {idx+1} — {loja}", use_column_width=True)

            col1, col2, col3 = st.columns([3, 1, 1])
            with col1:
                st.session_state.output_filename = st.text_input(
                    "Nome base do arquivo (sem extensão)",
                    value=st.session_state.get("output_filename", "Modelo_01"),
                    key="output_filename_input"
                )

            with col2:
                if st.session_state.get("ppt_bytes"):
                    st.download_button(
                        "⬇️ Baixar PPT",
                        data=st.session_state.ppt_bytes,
                        file_name=f"{(st.session_state.output_filename or 'Apresentacao')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentation.presentation",
                        use_container_width=True,
                        key=f"download_{st.session_state.get('download_key', 0)}"
                    )
                else:
                    btn_generate = st.button("🧩 Gerar PPT", key="btn_generate", use_container_width=True)

            with col3:
                if st.session_state.get("images_zip_bytes"):
                    st.download_button(
                        "⬇️ Baixar Imagens (ZIP)",
                        data=st.session_state.images_zip_bytes,
                        file_name=f"{(st.session_state.output_filename or 'Imagens')}.zip",
                        mime="application/zip",
                        use_container_width=True,
                        key=f"images_zip_{st.session_state.get('images_zip_key', 0)}"
                    )
                else:
                    btn_zip = st.button("🖼️ Baixar Imagens", key="btn_zip", use_container_width=True)

            if (not st.session_state.get("ppt_bytes")) and ('btn_generate' in locals()) and btn_generate:
                try:
                    titulo = (st.session_state.output_filename or "Apresentacao").strip()
                    use_template = cfg.get("use_template", False)
                    template_bytes = cfg.get("template_bytes")

                    if use_template and template_bytes:
                        ppt_bytes = gerar_ppt_modelo_capa_final(
                            template_bytes=template_bytes,
                            items=items, resultados=resultados, titulo=titulo,
                            max_per_slide=cfg["max_per_slide"], sort_mode=cfg["sort_mode"],
                            bg_rgb=cfg["bg_rgb"],
                            logo_bytes=cfg["logo_bytes"], logo_width_in=cfg["logo_width_in"],
                            signature_bytes=cfg["signature_bytes"],
                            signature_width_in=cfg.get("signature_width_in"),
                            auto_half_signature=cfg.get("auto_half_signature", True),
                            signature_bottom_margin_in=cfg["signature_bottom_margin_in"],
                            signature_right_margin_in=cfg["signature_right_margin_in"],
                            title_font_name=cfg["title_font_name"],
                            title_font_size_pt=cfg["title_font_size_pt"],
                            title_font_bold=cfg["title_font_bold"],
                            excluded_urls=st.session_state.excluded_urls
                        )
                    else:
                        prs = Presentation()
                        prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
                        blank = prs.slide_layouts[6]
                        title_rgb = pick_contrast_color(*cfg["bg_rgb"])
                        signature_width = (cfg["logo_width_in"]/2.0) if cfg.get("auto_half_signature", True) \
                                          else (cfg.get("signature_width_in") or 0.6)

                        groups = OrderedDict()
                        for loja, endereco, url in items:
                            if url in resultados and url not in st.session_state.excluded_urls:
                                groups.setdefault(str(loja), []).append((url, resultados[url]))

                        if cfg["sort_mode"] == "Nome da loja (A→Z)":
                            loja_keys = sorted(groups.keys(),
                                key=lambda s: (s is None or str(s).strip()== "", (s or "").strip().casefold()))
                        else:
                            loja_keys = list(groups.keys())

                        for loja in loja_keys:
                            imgs = groups[loja]; i = 0
                            while i < len(imgs):
                                if cfg["max_per_slide"] == "Automático":
                                    _url0, (_loja0, endereco, _buf0, (w0, h0), *_rest0) = imgs[i]
                                    per_slide = 3 if is_portrait(w0, h0) else 2
                                else:
                                    per_slide = int(cfg["max_per_slide"])
                                    endereco = imgs[i][1][1]

                                batch = imgs[i:i+per_slide]; i += per_slide
                                slide = prs.slides.add_slide(blank)
                                set_slide_bg(slide, cfg["bg_rgb"])
                                add_title_and_address(slide, loja, endereco, title_rgb,
                                    cfg["title_font_name"], cfg["title_font_size_pt"], cfg["title_font_bold"])
                                if cfg["logo_bytes"]:
                                    add_logo_top_right(slide, prs, cfg["logo_bytes"], cfg["logo_width_in"])
                                if cfg["signature_bytes"]:
                                    add_signature_bottom_right(slide, prs, cfg["signature_bytes"], signature_width,
                                        bottom_margin_in=cfg["signature_bottom_margin_in"],
                                        right_margin_in=cfg["signature_right_margin_in"])
                                slots = get_slots(len(batch), prs)
                                for (url, (_loja, _end, buf, (w_px, h_px), *rest)), (left, top, max_w_in, max_h_in) in zip(batch, slots):
                                    place_picture(slide, buf, w_px, h_px, left, top, max_w_in, max_h_in)

                        out = BytesIO(); prs.save(out); out.seek(0); ppt_bytes = out

                    st.session_state.ppt_bytes = ppt_bytes
                    st.session_state.generated = True
                    st.rerun()
                except Exception as e:
                    logger.exception("Falha ao gerar PPT")
                    st.error(f"Falha ao gerar PPT: {e}")

            if (not st.session_state.get("images_zip_bytes")) and ('btn_zip' in locals()) and btn_zip:
                try:
                    zip_bytes = montar_zip_imagens(
                        items=items,
                        resultados=resultados,
                        excluded_urls=st.session_state.excluded_urls
                    )
                    st.session_state.images_zip_bytes = zip_bytes
                    st.rerun()
                except Exception as e:
                    logger.exception("Falha ao montar ZIP")
                    st.error(f"Falha ao montar ZIP: {e}")

# -------------------------------------------------------------------
# ROTEAMENTO FINAL
# -------------------------------------------------------------------
if not st.session_state.auth:
    do_login()
else:
    main_app()
