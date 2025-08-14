# app.py
import re
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor, as_completed

import streamlit as st
import pandas as pd
import requests
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ============== CONFIG BÁSICA DA PÁGINA ==============
st.set_page_config(page_title="Gerador de Book", page_icon="📸", layout="wide")

# ============== TEMA (toggle Claro/Escuro) ==============
def apply_theme(dark: bool):
    # Isso "pinta" o app. Não troca o tema nativo do Streamlit, mas resolve visualmente.
    if dark:
        css = """
        <style>
        .stApp { background-color: #0e1117; color: #fafafa; }
        .stMarkdown, .stTextInput, .stFileUploader, .stButton, .stProgress { color: #fafafa; }
        .stSelectbox, .stTextInput > div > div > input { color: #fafafa; }
        </style>
        """
    else:
        css = """
        <style>
        .stApp { background-color: #ffffff; color: #1f2328; }
        </style>
        """
    st.markdown(css, unsafe_allow_html=True)

if "dark_mode" not in st.session_state:
    st.session_state.dark_mode = False

# ============== LOGIN SUPER SIMPLES (didático) ==============
ALLOWED_USERS = {
    "lucas.costa@mkthouse.com.br": "mudar12345",  # <<< TROQUE A SENHA AQUI SE QUISER
}

def do_login():
    st.title("🔐 Login")
    with st.form("login_form", clear_on_submit=False):
        email = st.text_input("E-mail", placeholder="seu.email@mkthouse.com.br")
        pwd = st.text_input("Senha", type="password", placeholder="••••••••")
        entrar = st.form_submit_button("Entrar")
    if entrar:
        if email in ALLOWED_USERS and pwd == ALLOWED_USERS[email]:
            st.session_state.auth = True
            st.session_state.user_email = email
            st.experimental_rerun()
        else:
            st.error("Credenciais inválidas. Verifique e tente novamente.")

if "auth" not in st.session_state:
    st.session_state.auth = False

# ============== FUNÇÕES DE IMAGEM E PPT ==============
URL_RE = re.compile(r'https?://\S+')

def extrair_links(celula):
    if pd.isna(celula):
        return []
    texto = str(celula).replace(",", " ")
    texto = texto.replace("(", " ").replace(")", " ").replace('"', " ").replace("'", " ")
    return [u.rstrip(").,") for u in URL_RE.findall(texto)]

def redimensionar(img: Image.Image, max_w: int, max_h: int) -> Image.Image:
    img = ImageOps.exif_transpose(img)  # corrige rotação
    if img.mode != "RGB":
        img = img.convert("RGB")
    img.thumbnail((max_w, max_h), resample=Image.LANCZOS)
    return img

def comprimir_jpeg_binsearch(img: Image.Image, limite_kb: int) -> BytesIO:
    lo, hi = 35, 95
    best_buf = None

    # tentativa rápida
    q_try = 75
    buf = BytesIO()
    img.save(buf, format="JPEG", quality=q_try, optimize=True, progressive=True, subsampling=2)
    if buf.tell() / 1024 <= limite_kb:
        buf.seek(0)
        return buf
    best_buf = buf

    while lo <= hi:
        mid = (lo + hi) // 2
        buf = BytesIO()
        img.save(buf, format="JPEG", quality=mid, optimize=True, progressive=True, subsampling=2)
        if buf.tell() / 1024 <= limite_kb:
            best_buf = buf
            lo = mid + 1
        else:
            hi = mid - 1

    if best_buf is None:
        best_buf = BytesIO()
        img.save(best_buf, format="JPEG", quality=35, optimize=True, progressive=True, subsampling=2)

    best_buf.seek(0)
    return best_buf

def baixar_processar(session, url: str, max_w: int, max_h: int, limite_kb: int, timeout: int):
    try:
        r = session.get(url, timeout=timeout, stream=True)
        if r.status_code != 200:
            return (url, False, None, None)
        img = Image.open(BytesIO(r.content))
        img = redimensionar(img, max_w, max_h)
        buf = comprimir_jpeg_binsearch(img, limite_kb)
        w, h = Image.open(buf).size
        buf.seek(0)
        return (url, True, buf, (w, h))
    except Exception:
        return (url, False, None, None)

def px_to_inches(px):
    return Inches(px / 96.0)  # aproximação 96 dpi

def gerar_ppt(items, resultados, titulo):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    TITLE_LEFT = Inches(0.5)
    TITLE_TOP = Inches(0.2)
    TITLE_W = Inches(12)
    TITLE_H = Inches(1)
    IMG_TOP = Inches(1.2)
    IMG_MAX_W = Inches(11)
    IMG_MAX_H = Inches(6)

    for loja, url in items:
        if url not in resultados:
            continue
        _, buf, (w_px, h_px) = resultados[url]
        slide = prs.slides.add_slide(blank)

        # título
        tx = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_W, TITLE_H)
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(loja)
        font = run.font
        font.name = 'Arial'
        font.size = Pt(15)
        font.bold = True
        font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = 1  # centralizado

        # imagem
        img_w_in = min(px_to_inches(w_px), IMG_MAX_W)
        img_h_in = min(px_to_inches(h_px), IMG_MAX_H)
        ratio = min(float(IMG_MAX_W) / float(img_w_in), float(IMG_MAX_H) / float(img_h_in), 1.0)
        final_w = img_w_in * ratio
        final_h = img_h_in * ratio
        img_left = (prs.slide_width - final_w) / 2
        img_top = IMG_TOP

        buf.seek(0)
        slide.shapes.add_picture(buf, img_left, img_top, width=final_w, height=final_h)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# ============== APP (fluxo principal) ==============
def main_app():
    # sidebar
    with st.sidebar:
        st.header("⚙️ Preferências")
        st.session_state.dark_mode = st.toggle("Usar tema escuro", value=st.session_state.dark_mode)
        apply_theme(st.session_state.dark_mode)

        st.markdown("---")
        st.caption("Colunas da planilha (nomes de cabeçalho):")
        loja_col = st.text_input("Coluna de LOJA", value="Selecione sua loja")
        img_col  = st.text_input("Coluna de FOTOS", value="Faça o upload das fotos")

        st.markdown("---")
        st.caption("Tamanho e compressão das imagens")
        target_w = st.number_input("Largura máx (px)",  min_value=480, max_value=4096, value=1280, step=10)
        target_h = st.number_input("Altura máx (px)",   min_value=360, max_value=4096, value=720, step=10)
        limite_kb = st.number_input("Tamanho máx por foto (KB)", min_value=50, max_value=2000, value=450, step=10)

        st.markdown("---")
        st.caption("Rede e paralelismo")
        max_workers = st.slider("Trabalhos em paralelo", min_value=2, max_value=32, value=12)
        req_timeout = st.slider("Timeout por download (s)", min_value=5, max_value=60, value=15)

    st.title("📸 Gerador de Book (PPT)")
    st.write("Arraste sua planilha Excel aqui (com os links das fotos).")

    up = st.file_uploader("Selecione ou arraste a planilha (.xlsx)", type=["xlsx"])
    gerar = st.button("🚀 Gerar PPT")

    if gerar and not up:
        st.warning("Envie a planilha primeiro.")
        st.stop()

    if up and gerar:
        # Lê planilha
        try:
            df = pd.read_excel(up)
        except Exception as e:
            st.error(f"Não consegui ler o Excel: {e}")
            st.stop()

        # Checa colunas
        missing = [c for c in [img_col, loja_col] if c not in df.columns]
        if missing:
            st.error(f"Colunas não encontradas: {missing}\nVerifique o cabeçalho da planilha.")
            st.stop()

        # Constroi lista (loja, url)
        items = []
        for _, row in df.iterrows():
            loja = str(row[loja_col]).strip()
            for url in extrair_links(row.get(img_col, "")):
                if url.startswith("http"):
                    items.append((loja, url))

        # Remove duplicados
        seen = set()
        uniq = []
        for loja, url in items:
            if url not in seen:
                seen.add(url)
                uniq.append((loja, url))
        items = uniq

        total = len(items)
        if total == 0:
            st.warning("Nenhuma URL de imagem encontrada na coluna de fotos.")
            st.stop()

        st.info(f"Serão processadas **{total}** imagens.")

        # Sessão HTTP otimizada
        session = requests.Session()
        adapter = requests.adapters.HTTPAdapter(pool_connections=max_workers, pool_maxsize=max_workers, max_retries=2)
        session.mount("http://", adapter)
        session.mount("https://", adapter)
        session.headers.update({"User-Agent": "Mozilla/5.0 (GeradorBook Streamlit)"})

        # Progresso
        prog = st.progress(0)
        status = st.empty()

        resultados = {}
        falhas = 0

        with ThreadPoolExecutor(max_workers=max_workers) as ex:
            futures = {ex.submit(baixar_processar, session, url, target_w, target_h, limite_kb, req_timeout): (loja, url)
                       for loja, url in items}
            done_count = 0
            for fut in as_completed(futures):
                loja, url = futures[fut]
                ok_url, ok, buf, wh = fut.result()
                if ok:
                    resultados[url] = (loja, buf, wh)
                else:
                    falhas += 1
                done_count += 1
                prog.progress(int(done_count * 100 / total))
                status.write(f"Processadas {done_count}/{total} imagens...")

        status.write(f"Processo concluído. Falhas: {falhas}")

        # Gera PPT em memória
        titulo = "Apresentacao_Relatorio_Compacta"
        ppt_bytes = gerar_ppt(items, resultados, titulo)
        st.success("PPT gerado com sucesso!")
        st.download_button(
            "⬇️ Baixar PPT",
            data=ppt_bytes,
            file_name=f"{titulo}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

# ============== ROTEAMENTO: login -> app ==============
if not st.session_state.auth:
    do_login()
else:
    # Topbar com usuário + sair
    topcol1, topcol2 = st.columns([1,1])
    with topcol1:
        st.caption(f"Logado como: **{st.session_state.user_email}**")
    with topcol2:
        if st.button("Sair", type="secondary"):
            st.session_state.clear()
            st.experimental_rerun()
    main_app()
